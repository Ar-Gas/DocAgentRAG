import os
import logging
import chardet
import re
import tempfile
import shutil
import subprocess
from pathlib import Path
from email import policy
from email.parser import BytesParser

from config import MAX_FILE_SIZE, MAX_TEXT_LENGTH, PDF_PAGE_LIMIT, EXCEL_CHUNK_SIZE

logger = logging.getLogger(__name__)

# ===================== 工具：图片处理（OCR提取文字）=====================
def process_image(filepath):
    """
    处理图片文件，使用OCR提取文字
    :param filepath: 图片路径
    :return: str
    """
    try:
        from PIL import Image
        import pytesseract
        
        logger.info(f"处理图片：{filepath}")
        
        # 打开图片
        img = Image.open(filepath)
        
        # 尝试OCR识别
        try:
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            if text and text.strip():
                content = text.strip()
                content = _truncate_text(content)
                logger.info(f"图片OCR识别成功，提取文本长度：{len(content)}")
                return content
            else:
                logger.warning(f"图片未识别到文字：{filepath}")
                return "图片中未识别到文字"
        except Exception as e:
            logger.error(f"OCR识别失败: {str(e)}")
            return f"图片OCR识别失败: {str(e)}"
            
    except Exception as e:
        logger.error(f"图片处理失败: {str(e)}")
        return f"图片处理失败: {str(e)}"

# ===================== 工具：文件校验 =====================
def _check_file_validity(filepath):
    """检查文件是否存在、大小是否合法"""
    if not os.path.exists(filepath):
        logger.error(f"文件不存在：{filepath}")
        return False, "文件不存在"
    if os.path.getsize(filepath) > MAX_FILE_SIZE:
        logger.error(f"文件过大（超过{MAX_FILE_SIZE//1024//1024}MB）：{filepath}")
        return False, f"文件过大（超过{MAX_FILE_SIZE//1024//1024}MB）"
    return True, ""

# ===================== 工具：文本截断 =====================
def _truncate_text(text):
    """截断超长文本，避免内存溢出"""
    if len(text) > MAX_TEXT_LENGTH:
        logger.warning(f"文本过长，截断至{MAX_TEXT_LENGTH}字符")
        return text[:MAX_TEXT_LENGTH] + "\n（文本过长，已截断）"
    return text

# ===================== 工具：扫描版PDF检测 =====================
def _is_scanned_pdf(filepath, min_text_length=100):
    """
    辅助函数：检测PDF是否为扫描版
    :param filepath: PDF路径
    :param min_text_length: 最小文本长度阈值（低于此值判定为扫描版）
    :return: bool
    """
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        total_text = ""
        # 仅检查前10页，避免大文件耗时过长
        check_pages = min(len(reader.pages), 10)
        for page in reader.pages[:check_pages]:
            text = page.extract_text()
            if text:
                total_text += text.strip()
        # 如果总文本长度小于阈值，判定为扫描版
        return len(total_text) < min_text_length
    except Exception as e:
        logger.warning(f"检测PDF类型失败，默认按扫描版处理: {str(e)}")
        return True

# ===================== 工具：MinerU OCR处理 =====================
def process_scanned_pdf_with_mineru(filepath):
    """
    使用MinerU处理扫描版PDF（OCR识别）
    :param filepath: PDF路径
    :return: (success: bool, content: str)
    """
    try:
        # 1. 检查MinerU依赖是否安装
        try:
            from magic_pdf.cli import magic_pdf
        except ImportError:
            logger.error("MinerU未安装，请执行：pip install magic-pdf[full]")
            return False, "MinerU未安装，请安装后重试（pip install magic-pdf[full]）"

        logger.info(f"检测到扫描版PDF，启动MinerU OCR处理：{filepath}")

        # 2. 创建临时目录（处理完自动清理）
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_pdf = temp_path / Path(filepath).name
            shutil.copy2(filepath, input_pdf)

            # 3. 调用MinerU处理（限制页数+轻量化模式）
            result = subprocess.run(
                [
                    "magic-pdf",
                    "-p", str(input_pdf),
                    "-o", str(temp_path),
                    "-m", "auto",
                    "--lang", "chinese+english"
                ],
                capture_output=True,
                text=True,
                timeout=600
            )

            if result.returncode != 0:
                logger.error(f"MinerU处理失败: {result.stderr}")
                return False, f"MinerU处理失败: {result.stderr[:500]}"

            # 4. 解析MinerU输出（提取Markdown文本）
            md_files = list(temp_path.glob("*.md"))
            if not md_files:
                logger.error("MinerU未生成输出文件")
                return False, "MinerU处理失败：未生成输出文件"

            # 读取Markdown内容
            with open(md_files[0], 'r', encoding='utf-8') as f:
                content = f.read()

            # 5. 轻量化处理（保留文本，去除复杂Markdown格式）
            content = re.sub(r'!\[.*?\]\(.*?\)', '', content)
            content = re.sub(r'\n{3,}', '\n\n', content)
            content = content.strip()

            # 6. 截断超长文本（防内存溢出）
            content = _truncate_text(content)

            logger.info(f"MinerU处理成功，提取文本长度：{len(content)}")
            return True, content

    except subprocess.TimeoutExpired:
        logger.error("MinerU处理超时（超过10分钟）")
        return False, "MinerU处理超时（文件过大或复杂度过高）"
    except Exception as e:
        logger.error(f"MinerU处理异常: {str(e)}")
        return False, f"MinerU处理异常: {str(e)}"

# ===================== PDF文档处理（普通+扫描版自动切换）=====================
def process_pdf(filepath):
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        content = []
        
        # 限制处理页数，防大文件
        total_pages = min(len(reader.pages), PDF_PAGE_LIMIT)
        logger.info(f"处理PDF：{filepath}，共{total_pages}页（限制{PDF_PAGE_LIMIT}页）")
        
        for page_num in range(total_pages):
            page = reader.pages[page_num]
            text = page.extract_text()
            if text and text.strip():
                content.append(f"--- 第 {page_num+1} 页 ---\n{text.strip()}")
        
        # 核心：检测到扫描版时，自动调用MinerU
        if not content or len('\n'.join(content)) < 100:
            logger.warning(f"PDF文本提取量极少，判定为扫描版，切换至MinerU OCR：{filepath}")
            success, ocr_content = process_scanned_pdf_with_mineru(filepath)
            if success:
                return ocr_content
            else:
                return f"（扫描版PDF，MinerU处理失败：{ocr_content}）"
        
        return _truncate_text('\n'.join(content))
    except Exception as e:
        logger.error(f"PDF处理失败: {str(e)}")
        return f"PDF处理失败: {str(e)}"

# ===================== Word文档处理（增加表格提取）=====================
def process_word(filepath):
    try:
        import docx
        doc = docx.Document(filepath)
        content = []
        logger.info(f"处理Word：{filepath}")
        
        # 提取段落
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                content.append(text)
        
        # 提取表格
        if doc.tables:
            content.append("\n--- 表格内容 ---")
            for table_idx, table in enumerate(doc.tables, 1):
                content.append(f"\n表格 {table_idx}:")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        content.append(row_text)
        
        return _truncate_text('\n'.join(content))
    except Exception as e:
        logger.error(f"Word处理失败: {str(e)}")
        return f"Word处理失败: {str(e)}"

# ===================== Excel文档处理（分块读取+优化内存）=====================
def process_excel(filepath):
    try:
        import pandas as pd
        content = []
        logger.info(f"处理Excel：{filepath}")
        
        # 读取所有工作表
        xls = pd.ExcelFile(filepath)
        for sheet_name in xls.sheet_names:
            content.append(f"\n--- 工作表: {sheet_name} ---")
            # 分块读取，防大文件
            for chunk_idx, df in enumerate(pd.read_excel(filepath, sheet_name=sheet_name, chunksize=EXCEL_CHUNK_SIZE)):
                if chunk_idx > 0:
                    content.append(f"\n（第 {chunk_idx+1} 块数据）")
                # 优化格式，避免超长
                content.append(df.to_string(index=False, na_rep='-', max_colwidth=20))
        
        return _truncate_text('\n'.join(content))
    except Exception as e:
        logger.error(f"Excel处理失败: {str(e)}")
        return f"Excel处理失败: {str(e)}"

# ===================== PPT文档处理 =====================
def process_ppt(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        content = []
        logger.info(f"处理PPT：{filepath}")
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_content = []
            # 提取文本框
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            # 提取表格
            if hasattr(slide.shapes, "tables"):
                for table in slide.shapes.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            slide_content.append(row_text)
            
            if slide_content:
                content.append(f"\n--- 第 {slide_idx} 页 ---\n" + "\n".join(slide_content))
        
        return _truncate_text('\n'.join(content))
    except Exception as e:
        logger.error(f"PPT处理失败: {str(e)}")
        return f"PPT处理失败: {str(e)}"

# ===================== 邮件处理（优化编码）=====================
def process_email(filepath):
    try:
        with open(filepath, 'rb') as f:
            msg = BytesParser(policy=policy.default).parse(f)
        content = []
        logger.info(f"处理邮件：{filepath}")
        
        # 提取头部信息
        if msg['From']:
            content.append(f"发件人: {msg['From']}")
        if msg['To']:
            content.append(f"收件人: {msg['To']}")
        if msg['Subject']:
            content.append(f"主题: {msg['Subject']}")
        if msg['Date']:
            content.append(f"发送时间: {msg['Date']}")
        
        # 提取正文
        content.append("\n--- 邮件正文 ---")
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    try:
                        body = part.get_content()
                        content.append(body)
                    except Exception:
                        pass
        else:
            try:
                body = msg.get_content()
                content.append(body)
            except Exception:
                pass
        
        return _truncate_text('\n'.join(content))
    except Exception as e:
        logger.error(f"邮件处理失败: {str(e)}")
        return f"邮件处理失败: {str(e)}"

# ===================== 统一文档处理接口（优化校验+错误返回）=====================
def process_document(filepath):
    """
    统一文档处理接口
    :param filepath: 文件路径
    :return: (success: bool, content: str)
    """
    # 先校验文件
    is_valid, error_msg = _check_file_validity(filepath)
    if not is_valid:
        return False, error_msg
    
    ext = os.path.splitext(filepath)[1].lower()
    logger.info(f"开始处理文件：{filepath}，类型：{ext}")
    
    # 处理器映射
    handlers = {
        '.pdf': process_pdf,
        '.docx': process_word,
        '.xlsx': process_excel,
        '.xls': process_excel,
        '.ppt': process_ppt,
        '.pptx': process_ppt,
        '.eml': process_email,
        '.msg': process_email,
        '.jpg': process_image,
        '.jpeg': process_image,
        '.png': process_image,
        '.gif': process_image,
        '.bmp': process_image,
        '.webp': process_image
    }
    
    # 调用对应处理器
    if ext in handlers:
        content = handlers[ext](filepath)
    else:
        # 处理其他文本文件（分块读取）
        try:
            content = []
            with open(filepath, 'rb') as f:
                while True:
                    chunk = f.read(1024 * 1024)  # 每次读1MB
                    if not chunk:
                        break
                    encoding = chardet.detect(chunk)['encoding'] or 'utf-8'
                    content.append(chunk.decode(encoding, errors='ignore'))
            content = '\n'.join(content)
            content = _truncate_text(content)
        except Exception as e:
            logger.error(f"文本文件处理失败: {str(e)}")
            return False, f"文本文件处理失败: {str(e)}"
    
    # 判断是否成功
    if content.startswith("处理失败") or content.startswith("文件过大") or content.startswith("文件不存在") or content.startswith("（扫描版PDF"):
        return False, content
    return True, content