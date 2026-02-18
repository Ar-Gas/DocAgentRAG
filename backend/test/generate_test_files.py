#!/usr/bin/env python3
import os
import sys
import subprocess
import random
import string
from datetime import datetime

# Check and install required dependencies
def install_dependencies():
    required_packages = [
        'python-docx',
        'pandas',
        'openpyxl',
        'fpdf',
        'python-pptx'
    ]
    
    print("Checking dependencies...")
    for package in required_packages:
        try:
            __import__(package.replace('-', '_'))
            print(f"✓ {package} is installed")
        except ImportError:
            print(f"Installing {package}...")
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', package])
            print(f"✓ {package} installed successfully")

# Generate random content
def generate_random_text(length=1000):
    words = [
        'Lorem', 'ipsum', 'dolor', 'sit', 'amet', 'consectetur', 'adipiscing', 'elit',
        'sed', 'do', 'eiusmod', 'tempor', 'incididunt', 'ut', 'labore', 'et', 'dolore',
        'magna', 'aliqua', 'Ut', 'enim', 'ad', 'minim', 'veniam', 'quis', 'nostrud',
        'exercitation', 'ullamco', 'laboris', 'nisi', 'ut', 'aliquip', 'ex', 'ea', 'commodo',
        'consequat', 'Duis', 'aute', 'irure', 'dolor', 'in', 'reprehenderit', 'in',
        'voluptate', 'velit', 'esse', 'cillum', 'dolore', 'eu', 'fugiat', 'nulla', 'pariatur',
        'Excepteur', 'sint', 'occaecat', 'cupidatat', 'non', 'proident', 'sunt', 'in',
        'culpa', 'qui', 'officia', 'deserunt', 'mollit', 'anim', 'id', 'est', 'laborum'
    ]
    text = []
    for _ in range(length):
        text.append(random.choice(words))
    return ' '.join(text)

# Generate PDF files
def generate_pdf_files(output_dir, count=10):
    from fpdf import FPDF
    
    print(f"Generating {count} PDF files...")
    for i in range(1, count + 1):
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add multiple pages
        for page_num in range(1, random.randint(3, 8)):
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            # Add title
            pdf.set_font("Arial", 'B', 16)
            pdf.cell(200, 10, txt=f"PDF Document {i} - Page {page_num}", ln=True, align='C')
            pdf.ln(10)
            
            # Add random text paragraphs
            pdf.set_font("Arial", size=12)
            for _ in range(random.randint(5, 10)):
                pdf.multi_cell(0, 10, txt=generate_random_text(random.randint(200, 500)))
                pdf.ln(5)
        
        file_path = os.path.join(output_dir, f"test_pdf_{i}.pdf")
        pdf.output(file_path)
        print(f"✓ Generated {file_path}")

# Generate PowerPoint files
def generate_pptx_files(output_dir, count=10):
    from pptx import Presentation
    from pptx.util import Inches
    
    print(f"Generating {count} PPT files...")
    for i in range(1, count + 1):
        prs = Presentation()
        
        # Add multiple slides
        for slide_num in range(1, random.randint(3, 8)):
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title = slide.shapes.title
            title.text = f"PPT Presentation {i} - Slide {slide_num}"
            
            # Add content
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.text = generate_random_text(random.randint(500, 1000))
        
        file_path = os.path.join(output_dir, f"test_ppt_{i}.pptx")
        prs.save(file_path)
        print(f"✓ Generated {file_path}")

# Generate Word files
def generate_docx_files(output_dir, count=10):
    from docx import Document
    
    print(f"Generating {count} Word files...")
    for i in range(1, count + 1):
        doc = Document()
        
        # Add title
        doc.add_heading(f"Word Document {i}", 0)
        
        # Add multiple sections
        for section_num in range(1, random.randint(3, 8)):
            doc.add_heading(f"Section {section_num}", level=1)
            
            # Add paragraphs
            for _ in range(random.randint(8, 15)):
                doc.add_paragraph(generate_random_text(random.randint(200, 400)))
        
        file_path = os.path.join(output_dir, f"test_word_{i}.docx")
        doc.save(file_path)
        print(f"✓ Generated {file_path}")

# Generate Email (.eml) files
def generate_eml_files(output_dir, count=10):
    print(f"Generating {count} Email files...")
    for i in range(1, count + 1):
        # Create eml content
        from_email = f"sender{i}@example.com"
        to_email = f"recipient{i}@example.com"
        subject = f"Test Email {i}"
        date = datetime.now().strftime('%a, %d %b %Y %H:%M:%S %z')
        
        # Generate body with multiple paragraphs
        body = "\n".join([generate_random_text(random.randint(300, 600)) for _ in range(random.randint(4, 8))])
        
        # Create eml structure
        eml_content = f"""From: {from_email}
To: {to_email}
Subject: {subject}
Date: {date}
MIME-Version: 1.0
Content-Type: text/plain; charset=utf-8

{body}
"""
        
        file_path = os.path.join(output_dir, f"test_email_{i}.eml")
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(eml_content)
        print(f"✓ Generated {file_path}")

# Generate Excel files
def generate_xlsx_files(output_dir, count=10):
    import pandas as pd
    
    print(f"Generating {count} Excel files...")
    for i in range(1, count + 1):
        # Create multiple sheets
        writer = pd.ExcelWriter(os.path.join(output_dir, f"test_excel_{i}.xlsx"), engine='openpyxl')
        
        for sheet_num in range(1, random.randint(3, 6)):
            # Create large dataframe
            rows = random.randint(1000, 3000)
            cols = random.randint(10, 20)
            
            data = {}
            for col in range(cols):
                col_name = f"Column_{chr(65 + col % 26)}{col // 26 + 1}"
                if random.random() < 0.5:
                    # Numeric data
                    data[col_name] = [random.uniform(1, 1000) for _ in range(rows)]
                else:
                    # Text data
                    data[col_name] = [generate_random_text(random.randint(10, 30)) for _ in range(rows)]
            
            df = pd.DataFrame(data)
            sheet_name = f"Sheet{sheet_num}"
            df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        writer.close()
        print(f"✓ Generated {output_dir}/test_excel_{i}.xlsx")

# Main function
def main():
    # Create output directory
    output_dir = "/root/autodl-tmp/DocAgentRAG/backend/test/test_date"
    os.makedirs(output_dir, exist_ok=True)
    
    print(f"Output directory: {output_dir}")
    
    # Install dependencies
    install_dependencies()
    
    # Generate files
    generate_pdf_files(output_dir)
    generate_pptx_files(output_dir)
    generate_docx_files(output_dir)
    generate_eml_files(output_dir)
    generate_xlsx_files(output_dir)
    
    print("\n✅ All test files generated successfully!")
    print(f"Total files: 50 (10 of each type)")
    print(f"Files saved to: {output_dir}")

if __name__ == "__main__":
    main()